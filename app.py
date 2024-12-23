import streamlit as st
import tempfile
from pathlib import Path
from pptx import Presentation
import csv
import re
import logging
from concurrent.futures import ThreadPoolExecutor

# Password Protection
PREDEFINED_PASSWORD = "securepassword123"

def password_protection():
    if "authenticated" not in st.session_state:
        st.session_state.authenticated = False
    if not st.session_state.authenticated:
        with st.form("password_form", clear_on_submit=True):
            password_input = st.text_input("Enter Password", type="password")
            submitted = st.form_submit_button("Submit")
            if submitted and password_input == PREDEFINED_PASSWORD:
                st.session_state.authenticated = True
                st.success("Access Granted! Please click 'Submit' again to proceed.")
            elif submitted:
                st.error("Incorrect Password")
        return False
    return True

# Decimal Consistency Validation
def validate_decimal_consistency(slide, slide_index):
    issues = []
    decimal_pattern = re.compile(r'\d+[\.,]\d+')  # Pattern to match decimal numbers with either '.' or ','
    decimal_places_set = set()

    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    matches = decimal_pattern.findall(run.text)
                    for match in matches:
                        logging.debug(f"Found decimal: {match}")  # Debugging line
                        decimal_places = len(match.split(',')[1] if ',' in match else match.split('.')[1])
                        decimal_places_set.add(decimal_places)

    if len(decimal_places_set) > 1:
        for match in matches:
            issues.append({
                'slide': slide_index,
                'issue': 'Inconsistent Decimal Points',
                'text': match,
                'details': f'Found inconsistent decimal points: {list(decimal_places_set)}'
            })
    return issues

# Save results to CSV
def save_to_csv(issues, output_csv):
    with open(output_csv, mode='w', newline='', encoding='utf-8') as file:
        writer = csv.DictWriter(file, fieldnames=['slide', 'issue', 'text', 'corrected', 'details'])
        writer.writeheader()
        writer.writerows(issues)

def main():
    if not password_protection():
        return

    st.title("Decimal Consistency Validator")
    uploaded_file = st.file_uploader("Upload a PowerPoint file", type=["pptx"])

    if uploaded_file:
        with tempfile.TemporaryDirectory() as tmpdir:
            temp_ppt_path = Path(tmpdir) / "uploaded_ppt.pptx"
            with open(temp_ppt_path, "wb") as f:
                f.write(uploaded_file.getbuffer())

            presentation = Presentation(temp_ppt_path)
            total_slides = len(presentation.slides)

            # Run Validation
            if st.button("Run Decimal Validation"):
                progress_bar = st.progress(0)
                progress_text = st.empty()
                issues = []

                # Parallel Processing
                with ThreadPoolExecutor() as executor:
                    futures = []
                    for slide_index in range(total_slides):
                        slide = presentation.slides[slide_index]
                        futures.append(executor.submit(validate_decimal_consistency, slide, slide_index + 1))

                    for i, future in enumerate(futures):
                        issues.extend(future.result())
                        progress_percent = int((i + 1) / len(futures) * 100)
                        progress_text.text(f"Progress: {progress_percent}%")
                        progress_bar.progress(progress_percent / 100)

                # Save Results
                csv_output_path = Path(tmpdir) / "decimal_validation_report.csv"
                save_to_csv(issues, csv_output_path)

                # Store results in session state
                st.session_state['csv_output'] = csv_output_path.read_bytes()
                st.success("Decimal validation completed!")

                # Display Download Buttons
                if 'csv_output' in st.session_state:
                    st.download_button("Download Validation Report (CSV)", st.session_state['csv_output'], file_name="decimal_validation_report.csv")

if __name__ == "__main__":
    main()